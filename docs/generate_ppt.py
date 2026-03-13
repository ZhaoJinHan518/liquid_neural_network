"""
SITP 结题汇报 PPT 生成脚本
项目：液态神经网络与普通神经网络的对比研究

用法：
  python docs/generate_ppt.py                        # 生成 PPTX（默认路径）
  python docs/generate_ppt.py --out my_slides.pptx  # 自定义输出路径
  python docs/generate_ppt.py --pdf                  # 同时导出 PDF（需要 LibreOffice）
  python docs/generate_ppt.py --pdf --out /tmp/sitp.pptx  # 自定义输出路径，同时生成 .pptx 和 .pdf
"""

import argparse
import os
import subprocess
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 颜色主题 ──────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0D, 0x1B, 0x2A)   # 深蓝背景
C_ACCENT   = RGBColor(0x00, 0xB4, 0xD8)   # 青蓝强调色
C_ACCENT2  = RGBColor(0x90, 0xE0, 0xEF)   # 浅青
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT    = RGBColor(0xCA, 0xE9, 0xFF)
C_GRAY     = RGBColor(0xA8, 0xB2, 0xC1)
C_YELLOW   = RGBColor(0xFF, 0xD1, 0x66)
C_GREEN    = RGBColor(0x52, 0xD4, 0x8B)
C_RED      = RGBColor(0xFF, 0x6B, 0x6B)
C_HEADER   = RGBColor(0x02, 0x3E, 0x5E)   # 深蓝标题栏

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── 工具函数 ──────────────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = line_width
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True, font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline_text(slide, lines, left, top, width, height,
                       font_size=16, bold=False, color=C_WHITE,
                       align=PP_ALIGN.LEFT, line_spacing=1.2, font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        # Support (text, bold, color, size) tuples or plain strings
        if isinstance(line, tuple):
            txt, b, c, s = (line + (None, None, None))[:4]
            run.text = txt
            run.font.bold = b if b is not None else bold
            run.font.color.rgb = c if c is not None else color
            run.font.size = Pt(s) if s is not None else Pt(font_size)
        else:
            run.text = line
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.size = Pt(font_size)
        run.font.name = font_name
    return txBox


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_divider(slide, top, color=C_ACCENT, width_in=12.0, left=0.67):
    bar = add_rect(slide, left, top, width_in, 0.04, fill_color=color)
    return bar


def slide_header(slide, title, subtitle=None):
    """Gradient-like top banner + title."""
    add_rect(slide, 0, 0, 13.33, 1.3, fill_color=C_HEADER)
    add_rect(slide, 0, 1.3, 13.33, 0.05, fill_color=C_ACCENT)
    add_text(slide, title, 0.5, 0.15, 12.0, 0.8,
             font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.85, 12.0, 0.4,
                 font_size=14, color=C_ACCENT2, align=PP_ALIGN.LEFT)


def make_table(slide, headers, rows,
               left, top, col_widths, row_height=0.42,
               header_fill=C_HEADER, header_color=C_WHITE,
               row_fill=RGBColor(0x06, 0x2A, 0x40),
               alt_fill=RGBColor(0x04, 0x20, 0x35),
               text_color=C_WHITE, font_size=13, bold_header=True):
    """Draw a simple table using rectangles and text boxes."""
    n_rows = len(rows)
    n_cols = len(headers)
    # Header
    x = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, top, w, row_height, fill_color=header_fill)
        add_text(slide, h, x + 0.05, top + 0.05, w - 0.1, row_height - 0.1,
                 font_size=font_size, bold=bold_header, color=header_color, align=PP_ALIGN.CENTER)
        x += w
    # Rows
    for i, row in enumerate(rows):
        y = top + row_height * (i + 1)
        x = left
        fill = row_fill if i % 2 == 0 else alt_fill
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_rect(slide, x, y, w, row_height, fill_color=fill)
            cell_color = text_color
            cell_bold = False
            if isinstance(cell, tuple):
                cell_text, cell_color, cell_bold = (cell + (text_color, False))[:3]
            else:
                cell_text = cell
            add_text(slide, cell_text, x + 0.05, y + 0.05, w - 0.1, row_height - 0.1,
                     font_size=font_size, bold=cell_bold, color=cell_color, align=PP_ALIGN.CENTER)
            x += w


def bullet_list(slide, items, left, top, width, height,
                font_size=16, color=C_WHITE, bullet="◆", indent=0.3, spacing=0.55, font_name="微软雅黑"):
    y = top
    for item in items:
        if isinstance(item, tuple):
            bul, txt, fs, c, b = (item + (font_size, color, False))[:5]
        else:
            bul, txt, fs, c, b = bullet, item, font_size, color, False
        if bul:
            add_text(slide, bul, left, y, indent, spacing, font_size=fs, bold=True, color=C_ACCENT)
            add_text(slide, txt, left + indent, y, width - indent, spacing,
                     font_size=fs, color=c, bold=b, font_name=font_name)
        else:
            add_text(slide, txt, left, y, width, spacing,
                     font_size=fs, color=c, bold=b, font_name=font_name)
        y += spacing


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, C_BG)

    # Top decorative bar
    add_rect(slide, 0, 0, 13.33, 0.12, fill_color=C_ACCENT)
    # Bottom bar
    add_rect(slide, 0, 7.38, 13.33, 0.12, fill_color=C_ACCENT)

    # Side accent stripe
    add_rect(slide, 0, 0.12, 0.12, 7.26, fill_color=C_ACCENT2)

    # Central content box
    add_rect(slide, 0.5, 0.8, 12.3, 5.5,
             fill_color=RGBColor(0x06, 0x2E, 0x4E))

    # Title
    add_text(slide, "液态神经网络与普通神经网络的对比研究",
             0.7, 1.1, 11.9, 1.0,
             font_size=36, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    add_divider(slide, 2.25, width_in=11.0, left=1.15)

    # Subtitle
    add_text(slide, "Comparative Study of Liquid Neural Networks and Conventional Neural Networks",
             0.7, 2.4, 11.9, 0.6,
             font_size=16, italic=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

    # Tag
    add_rect(slide, 4.7, 3.15, 3.9, 0.45,
             fill_color=C_ACCENT, line_color=None)
    add_text(slide, "同济大学大学生创新训练项目（SITP）",
             4.75, 3.18, 3.8, 0.38,
             font_size=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    # Info
    info_lines = [
        "团队成员：陈山璞  韩兆津",
        "指导单位：同济大学",
        "研究周期：2025 年 4 月 — 2026 年 4 月",
        "汇报日期：2026 年 3 月",
    ]
    y = 3.85
    for line in info_lines:
        add_text(slide, line, 2.5, y, 8.3, 0.45,
                 font_size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)
        y += 0.48


def slide_02_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "目录", "Agenda")

    sections = [
        ("01", "项目背景与意义"),
        ("02", "研究目标与创新点"),
        ("03", "技术方案：液态神经网络原理"),
        ("04", "实验设计（7 项实验）"),
        ("05", "核心实验结果"),
        ("06", "综合分析与核心发现"),
        ("07", "项目成果与贡献"),
        ("08", "挑战与未来展望"),
        ("09", "参考文献"),
    ]

    cols = 3
    per_col = 3
    col_w = 4.1
    for idx, (num, title) in enumerate(sections):
        col = idx // per_col
        row = idx % per_col
        x = 0.4 + col * col_w
        y = 1.7 + row * 1.65

        add_rect(slide, x, y, 3.75, 1.35,
                 fill_color=RGBColor(0x04, 0x28, 0x45))
        add_rect(slide, x, y, 0.6, 1.35, fill_color=C_ACCENT)
        add_text(slide, num, x + 0.05, y + 0.35, 0.5, 0.6,
                 font_size=20, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.7, y + 0.35, 3.0, 0.65,
                 font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)


def slide_03_background(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "01  项目背景与意义", "Background & Motivation")

    # Left column: problems
    add_rect(slide, 0.4, 1.55, 5.8, 5.6,
             fill_color=RGBColor(0x06, 0x2E, 0x4E))
    add_text(slide, "传统 RNN 的局限性", 0.5, 1.65, 5.6, 0.5,
             font_size=18, bold=True, color=C_ACCENT)

    problems = [
        ("①", "刚性动力学问题", "现实系统中存在时间尺度相差悬殊的多维动力学（如校园下课时人流激增），"
               "传统离散时间更新机制无法精确捕捉极短时窗内的急剧跳变。"),
        ("②", "不规则采样问题", "传感器故障、网络丢包导致时间步不均匀，传统 RNN 等时间间隔假设在实际场景中频繁失效。"),
        ("③", "参数效率问题", "传统 RNN 参数量随任务复杂度快速增长，在资源受限的边缘计算场景中难以部署。"),
    ]
    y = 2.25
    for num, title, desc in problems:
        add_text(slide, num, 0.5, y, 0.4, 0.35, font_size=16, bold=True, color=C_YELLOW)
        add_text(slide, title, 0.9, y, 5.1, 0.35, font_size=15, bold=True, color=C_ACCENT2)
        add_text(slide, desc, 0.5, y + 0.38, 5.5, 0.9, font_size=12, color=C_LIGHT)
        y += 1.45

    # Right column: LNN solution
    add_rect(slide, 6.6, 1.55, 6.3, 5.6,
             fill_color=RGBColor(0x02, 0x38, 0x50))
    add_text(slide, "液态神经网络（LNN）的突破", 6.7, 1.65, 6.1, 0.5,
             font_size=18, bold=True, color=C_ACCENT)

    add_text(slide, "LTC（液态时间常数网络）", 6.7, 2.3, 6.0, 0.38,
             font_size=16, bold=True, color=C_GREEN)
    ltc_desc = ("Hasani 等（AAAI 2021）提出，引入基于常微分方程（ODE）的连续时间动力学，"
                "从数学层面原生支持不规则采样和刚性系统建模。时间常数 τ 随输入信号自适应调整。")
    add_text(slide, ltc_desc, 6.7, 2.75, 6.0, 1.1, font_size=13, color=C_LIGHT)

    add_divider(slide, 3.98, color=C_ACCENT2, width_in=5.8, left=6.65)

    add_text(slide, "CfC（闭合形式连续时间网络）", 6.7, 4.1, 6.0, 0.38,
             font_size=16, bold=True, color=C_GREEN)
    cfc_desc = ("Hasani 等（Nature Machine Intelligence 2022）提出，通过解析求解 ODE 避免了数值迭代开销，"
                "将推理延迟从 ~20 ms（LTC）降至 ~5.8 ms，同时保留连续时间动力学特性。")
    add_text(slide, cfc_desc, 6.7, 4.55, 6.0, 1.1, font_size=13, color=C_LIGHT)

    add_text(slide, "🔬 本项目目标：系统量化对比 LNN 与传统 RNN 的性能差异",
             6.7, 5.8, 6.0, 0.7, font_size=14, bold=True, color=C_YELLOW)


def slide_04_goals(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "02  研究目标与创新点", "Objectives & Innovations")

    # Left: Objectives
    add_rect(slide, 0.4, 1.55, 6.0, 5.6, fill_color=RGBColor(0x06, 0x2A, 0x40))
    add_text(slide, "研究目标", 0.55, 1.65, 5.7, 0.45, font_size=20, bold=True, color=C_ACCENT)

    goals = [
        ("1.", "理论梳理", "系统整理 LTC、CfC 的 ODE 数学原理，与 LSTM/GRU/RNN 门控机制进行理论对比。"),
        ("2.", "代码实现", "基于 PyTorch + ncps 库，构建五种模型的统一接口（src/models/registry.py）。"),
        ("3.", "多任务实验", "覆盖七类任务：时序预测 × 4、序列分类、鲁棒性评估、强化学习控制。"),
        ("4.", "创新应用", "将 CfC 首次应用于同济大学校园人流预测，验证刚性/缺失场景的优越性。"),
        ("5.", "开源贡献", "提供完整中文文档 + 可复现代码，降低国内研究者学习门槛。"),
    ]
    y = 2.2
    for num, title, desc in goals:
        add_text(slide, num, 0.5, y, 0.4, 0.38, font_size=15, bold=True, color=C_ACCENT2)
        add_text(slide, title, 0.9, y, 5.3, 0.38, font_size=15, bold=True, color=C_WHITE)
        add_text(slide, desc, 0.5, y + 0.38, 5.8, 0.72, font_size=12, color=C_GRAY)
        y += 1.1

    # Right: Innovations
    add_rect(slide, 6.8, 1.55, 6.1, 5.6, fill_color=RGBColor(0x04, 0x22, 0x38))
    add_text(slide, "创新点", 6.95, 1.65, 5.8, 0.45, font_size=20, bold=True, color=C_ACCENT)

    innovations = [
        ("★", "系统性横向对比框架",
         "5 种模型 × 7 类实验任务，国内首个中文可复现液态网络对比基准。"),
        ("★", "校园人流预测创新场景",
         "含 6 个刚性尖峰（σ=3 min）+ 15% 随机缺失的合成校园人流数据集。"),
        ("★", "多维鲁棒性定量评估",
         "噪声注入（σ∈{0.00, 0.01, 0.05}）+ 推理延迟（CPU×200次）+ R² 综合画像。"),
        ("★", "全链路工程规范",
         "固定种子(42)、超参数日志、一键运行，实验结果完全可复现。"),
    ]
    y = 2.2
    for star, title, desc in innovations:
        add_text(slide, star, 6.9, y, 0.4, 0.38, font_size=16, bold=True, color=C_YELLOW)
        add_text(slide, title, 7.3, y, 5.4, 0.38, font_size=15, bold=True, color=C_WHITE)
        add_text(slide, desc, 6.9, y + 0.4, 5.9, 0.78, font_size=12, color=C_GRAY)
        y += 1.3


def slide_05_models(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "03  技术方案：液态神经网络原理", "LNN Architecture & Implementation")

    # ODE equation box
    add_rect(slide, 0.4, 1.6, 6.0, 2.3, fill_color=RGBColor(0x03, 0x26, 0x40))
    add_text(slide, "LTC 状态方程（ODE）", 0.55, 1.7, 5.7, 0.45,
             font_size=17, bold=True, color=C_ACCENT)
    add_text(slide, "dx/dt  =  −[ 1/τ + f(x, I, θ) ] · x  +  f(x, I, θ)",
             0.55, 2.25, 5.75, 0.6,
             font_size=17, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
    add_text(slide,
             "• x：神经元状态     • τ：自适应时间常数\n"
             "• f(x, I, θ)：输入驱动函数\n"
             "• 信号剧变时 τ 自动缩短 → 快速响应\n"
             "• 信号平缓时 τ 自动延长 → 惯性演化",
             0.55, 2.9, 5.75, 0.95, font_size=13, color=C_LIGHT)

    # CfC box
    add_rect(slide, 0.4, 4.0, 6.0, 2.2, fill_color=RGBColor(0x02, 0x30, 0x45))
    add_text(slide, "CfC 闭合形式近似（解析解）", 0.55, 4.1, 5.7, 0.45,
             font_size=17, bold=True, color=C_GREEN)
    add_text(slide, "避免 ODE 数值迭代 → 推理速度提升约 3.4×",
             0.55, 4.62, 5.7, 0.38, font_size=14, bold=True, color=C_YELLOW)
    add_text(slide,
             "• 保留连续时间动力学核心特性\n"
             "• 推理延迟：LTC ~20 ms → CfC ~5.8 ms\n"
             "• 反向传播更稳定，分类任务表现更优",
             0.55, 5.08, 5.7, 1.0, font_size=13, color=C_LIGHT)

    # Right: Model comparison table
    add_text(slide, "五种模型统一接口对比", 6.7, 1.65, 6.3, 0.45,
             font_size=18, bold=True, color=C_ACCENT)

    headers = ["模型", "类型", "核心特性", "参数量*"]
    rows = [
        ["LTC",  "液态 LNN", "ODE 连续时间动力学", "5,542"],
        ["CfC",  "液态 LNN", "闭合形式 ODE，推理高效", "21,025"],
        ["LSTM", "传统 RNN", "门控记忆，梯度稳定", "4,641"],
        ["GRU",  "传统 RNN", "简化门控，参数较少", "3,489"],
        ["RNN",  "传统 RNN", "基线，无门控机制", "~2,600"],
    ]
    make_table(slide, headers, rows,
               left=6.7, top=2.2,
               col_widths=[1.0, 1.3, 2.7, 1.2],
               font_size=12, row_height=0.48)

    add_text(slide, "* 鲁棒性评估实验（32 隐藏单元，回归任务）",
             6.7, 5.2, 6.1, 0.35, font_size=11, italic=True, color=C_GRAY)

    # Code snippet
    add_rect(slide, 6.7, 5.6, 6.1, 1.55, fill_color=RGBColor(0x01, 0x18, 0x28))
    add_text(slide, "# 统一调用方式（registry.py）",
             6.8, 5.7, 6.0, 0.35, font_size=11, italic=True, color=C_GRAY)
    add_text(slide, 'model = build_model("cfc", input_size=1, units=32, output_size=1)',
             6.8, 6.05, 5.9, 0.35, font_size=12, color=C_ACCENT2, font_name="Courier New")
    add_text(slide, 'output, hidden = model(x, hx)',
             6.8, 6.42, 5.9, 0.35, font_size=12, color=C_ACCENT2, font_name="Courier New")
    add_text(slide, '# x: (batch, seq_len, input_size)',
             6.8, 6.78, 5.9, 0.35, font_size=11, italic=True, color=C_GRAY, font_name="Courier New")


def slide_06_exp_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "04  实验设计总览", "7 Experiments at a Glance")

    headers = ["#", "实验任务", "数据集", "对比模型", "核心指标"]
    rows = [
        ["1", "阻尼正弦波预测\n（论文 Fig 6 复现）",
         "合成正弦，1000 条，长度 100", "LTC vs GRU", "MSE / 训练速度"],
        ["2", "Walker2d 轨迹预测\n（Section 4.1 复现）",
         "合成/Minari，4800 条，17 维", "LTC vs LSTM", "MSE / 训练速度"],
        ["3", "校园人流预测 ★\n（SITP 创新模块）",
         "同济合成，60 天，15% 缺失", "CfC vs LSTM", "MSE / 5-step 滚动"],
        ["4", "多维鲁棒性评估",
         "合成正弦，噪声注入", "CfC/LTC/LSTM/GRU", "RMSE@σ / R² / 延迟"],
        ["5", "金融时序预测",
         "NASDAQ 对数收益率，5 年", "CfC/LSTM/GRU", "MSE / MAE"],
        ["6", "UCI HAR 分类",
         "9 维传感器，6 类活动，128 步", "LTC/CfC/LSTM/GRU/RNN", "准确率 / 参数量"],
        ["7", "CartPole RL 控制",
         "Gymnasium CartPole-v1", "CfC vs LSTM", "最终奖励 / 收敛速度"],
    ]
    make_table(slide, headers, rows,
               left=0.3, top=1.6,
               col_widths=[0.4, 3.0, 3.0, 2.8, 3.5],
               font_size=11, row_height=0.7)

    add_text(slide, "★ 创新应用场景", 0.35, 6.8, 3.0, 0.35,
             font_size=12, color=C_YELLOW, bold=True)


def slide_07_exp1_sine(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "05a  实验一：阻尼正弦波预测", "Exp 1 – Damped Sine Wave (Paper Fig 6 Reproduction)")

    headers = ["模型", "隐藏单元", "测试 MSE", "参数量", "训练时间/轮"]
    rows = [
        ["LTC", "32", "1.496×10⁻⁴", "5,380", "10.28 s"],
        [("GRU", C_GREEN, True), "32", ("1.287×10⁻⁴ ✓", C_GREEN, True), "3,393", ("0.57 s ✓", C_GREEN, True)],
    ]
    make_table(slide, headers, rows,
               left=0.4, top=1.6,
               col_widths=[1.5, 1.8, 2.5, 2.0, 2.5],
               font_size=14, row_height=0.55)

    findings = [
        "GRU（MSE=1.287×10⁻⁴）略优于 LTC（MSE=1.496×10⁻⁴）",
        "合成正弦波为规则平滑信号，无刚性跳变 → GRU 门控机制足以高效建模",
        "LTC 训练时间约为 GRU 的 18 倍（ODE 求解开销）",
        "两模型均可捕捉衰减趋势；LTC 在衰减尾部（低幅段）拟合偏差略大",
        "结论：液态网络优势在规则信号任务上尚不显著，核心价值体现于刚性/缺失场景",
    ]
    add_text(slide, "关键发现", 0.4, 3.05, 12.5, 0.45, font_size=17, bold=True, color=C_ACCENT)
    bullet_list(slide, findings, 0.4, 3.55, 12.5, 4.0, font_size=14, spacing=0.58)


def slide_08_exp2_walker(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "05b  实验二：Walker2d 轨迹预测", "Exp 2 – Walker2d Trajectory (Paper Section 4.1 Reproduction)")

    headers = ["模型", "隐藏单元", "测试 MSE", "参数量", "训练时间/轮"]
    rows = [
        ["LTC", "64", "0.08773", "26,180", "5.39 s"],
        [("LSTM", C_GREEN, True), "64", ("0.02036 ✓", C_GREEN, True), "22,353", ("0.20 s ✓", C_GREEN, True)],
    ]
    make_table(slide, headers, rows,
               left=0.4, top=1.6,
               col_widths=[1.5, 1.8, 2.5, 2.0, 2.5],
               font_size=14, row_height=0.55)

    findings = [
        "LSTM（MSE=0.02036）显著优于 LTC（MSE=0.08773），差距约 4.3×",
        "合成数据：均匀采样、无刚性跳变 → LSTM 离散门控具有优势",
        "LTC 训练时间约为 LSTM 的 27 倍（ODE 数值迭代开销）",
        "结论：液态网络优势主要体现于刚性动力学和不规则采样场景，而非规则合成数据",
    ]
    add_text(slide, "关键发现", 0.4, 3.05, 12.5, 0.45, font_size=17, bold=True, color=C_ACCENT)
    bullet_list(slide, findings, 0.4, 3.55, 12.5, 4.0, font_size=14, spacing=0.65)


def slide_09_campus_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "05c  实验三：校园人流预测（SITP 创新模块）★",
                 "Exp 3 – Campus Pedestrian Flow Prediction (SITP Innovation)")

    # Dataset description
    add_rect(slide, 0.4, 1.6, 5.8, 2.8, fill_color=RGBColor(0x04, 0x28, 0x42))
    add_text(slide, "数据集设计（同济合成）", 0.55, 1.7, 5.5, 0.42,
             font_size=16, bold=True, color=C_ACCENT)
    ds_lines = [
        "• 模拟 60 天 × 5 分钟间隔人流数据",
        "• 6 个下课时间点高斯尖峰（σ=3 min，峰值 ×1.5）",
        "  → 10:00 / 11:45 / 13:30 / 15:15 / 17:00 / 18:30",
        "• 15% 随机时间步缺失（掩码输入 [flow, mask]）",
        "• 高斯观测噪声（σ=0.02）",
        "• 序列长度 48 步（≈ 4 小时），训练 100 轮",
    ]
    y = 2.2
    for line in ds_lines:
        add_text(slide, line, 0.5, y, 5.6, 0.38, font_size=12, color=C_LIGHT)
        y += 0.38

    # Results table
    headers = ["模型", "缺失率", "测试 MSE", "参数量", "训练时间/轮"]
    rows = [
        [("CfC ★", C_GREEN, True), "15%",
         ("7.166×10⁻⁴ ✓", C_GREEN, True), "21,025", "11.58 s"],
        ["LSTM", "15%", "1.628×10⁻³", "4,641", "1.21 s"],
    ]
    make_table(slide, headers, rows,
               left=0.4, top=4.55,
               col_widths=[1.5, 1.5, 2.5, 2.0, 2.3],
               font_size=13, row_height=0.5)

    # Right: Analysis
    add_rect(slide, 6.5, 1.6, 6.4, 5.55, fill_color=RGBColor(0x02, 0x30, 0x45))
    add_text(slide, "为什么 CfC 胜出？", 6.65, 1.7, 6.1, 0.42,
             font_size=18, bold=True, color=C_ACCENT)

    reasons = [
        ("①", "刚性动力学建模",
         "CfC 连续时间 ODE 通过自适应 τ 精确捕捉下课铃引发的瞬间人流激增；"
         "LSTM 离散门控对每时间步一视同仁，难以区分平稳段与突变段。"),
        ("②", "缺失数据处理",
         "τ 在数据缺失（mask=0）时自动延长，状态以惯性方式平滑演化，"
         "实现隐式插值；LSTM 仅将掩码作为普通输入特征。"),
        ("③", "5-step 滚动预测",
         "自回归递推中 CfC 误差积累速度明显慢于 LSTM，"
         "连续时间动力学数值更稳定。"),
    ]
    y = 2.25
    for num, title, desc in reasons:
        add_text(slide, num, 6.6, y, 0.45, 0.38, font_size=16, bold=True, color=C_YELLOW)
        add_text(slide, title, 7.05, y, 5.7, 0.38, font_size=15, bold=True, color=C_WHITE)
        add_text(slide, desc, 6.6, y + 0.42, 6.2, 0.88, font_size=12, color=C_LIGHT)
        y += 1.42

    add_rect(slide, 6.5, 6.1, 6.4, 0.65, fill_color=RGBColor(0x00, 0x60, 0x80))
    add_text(slide, "CfC 的测试 MSE 仅为 LSTM 的 44%，性能提升超过一倍",
             6.6, 6.17, 6.2, 0.5, font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_10_robustness(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "05d  实验四：多维鲁棒性评估", "Exp 4 – Multi-Dimensional Robustness Evaluation")

    # Full results table
    headers = ["模型", "参数量", "推理延迟 (ms)", "RMSE (σ=0.00)", "RMSE (σ=0.01)", "RMSE (σ=0.05)", "R² 评分"]
    rows = [
        [("CfC", C_GREEN, True), "21,025", "5.775",
         ("0.02900 ✓", C_GREEN, True), ("0.03049 ✓", C_GREEN, True), "0.06422", ("0.9669 ✓", C_GREEN, True)],
        ["LTC", "5,542",  "19.791", "0.03394", "0.03507", "0.06080", "0.9547"],
        [("GRU", C_ACCENT2, False), "3,489",  ("0.825 ✓", C_GREEN, True),
         "0.03722", "0.03818", ("0.05729 ✓", C_GREEN, True), "0.9455"],
        ["LSTM", "4,641",  ("0.261 ✓", C_GREEN, True),
         "0.03994", "0.04091", ("0.05828 ✓", C_GREEN, True), "0.9372"],
    ]
    make_table(slide, headers, rows,
               left=0.3, top=1.6,
               col_widths=[1.2, 1.4, 2.0, 2.0, 2.0, 2.0, 1.9],
               font_size=12, row_height=0.52)

    # Findings
    findings = [
        ("准确性（洁净）", "CfC（RMSE=0.029，R²=0.967）> LTC > GRU > LSTM —— 液态网络整体优于传统 RNN"),
        ("高噪声鲁棒性", "σ=0.05 时，GRU/LSTM 误差增幅反而低于 CfC/LTC —— 门控低通滤波效果；液态网络优势在刚性跳变而非纯随机噪声"),
        ("推理延迟", "LSTM（0.261 ms）< GRU（0.825 ms）< CfC（5.775 ms）< LTC（19.791 ms）—— CfC 是液态网络中延迟最优选择"),
    ]
    add_text(slide, "三维分析", 0.3, 4.42, 12.7, 0.42, font_size=17, bold=True, color=C_ACCENT)
    y = 4.92
    for dim, content in findings:
        add_rect(slide, 0.3, y, 2.0, 0.6, fill_color=RGBColor(0x00, 0x5F, 0x80))
        add_text(slide, dim, 0.35, y + 0.1, 1.9, 0.42, font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, content, 2.4, y + 0.05, 10.5, 0.52, font_size=12, color=C_LIGHT)
        y += 0.72


def slide_11_other_exps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "05e  实验五～七：金融预测 / UCI HAR / CartPole RL",
                 "Exp 5–7 – Finance / HAR Classification / Reinforcement Learning")

    # Finance
    add_rect(slide, 0.3, 1.6, 4.15, 2.8, fill_color=RGBColor(0x04, 0x26, 0x3E))
    add_text(slide, "实验五：金融时序预测", 0.4, 1.7, 4.0, 0.4, font_size=15, bold=True, color=C_ACCENT)
    make_table(slide, ["模型", "测试 MSE", "测试 MAE"],
               [["CfC", "1.0833", "0.8319"],
                ["LSTM", "1.0783", "0.8309"],
                [("GRU ✓", C_GREEN, True), ("1.0778", C_GREEN, True), ("0.8307", C_GREEN, True)]],
               left=0.35, top=2.18, col_widths=[1.2, 1.5, 1.3],
               font_size=11, row_height=0.42)
    add_text(slide,
             "三种模型性能高度接近（Δ<0.006）\n"
             "金融序列接近随机游走 → LNN 额外代价\n"
             "难以转化为性能增益",
             0.4, 3.35, 4.0, 1.0, font_size=12, color=C_LIGHT)

    # UCI HAR
    add_rect(slide, 4.6, 1.6, 4.6, 2.8, fill_color=RGBColor(0x03, 0x24, 0x3C))
    add_text(slide, "实验六：UCI HAR 分类", 4.7, 1.7, 4.4, 0.4, font_size=15, bold=True, color=C_ACCENT)
    make_table(slide, ["模型", "准确率", "参数量"],
               [[("CfC", C_GREEN, True), ("100.0% ✓", C_GREEN, True), "23,142"],
                [("LSTM", C_GREEN, True), ("100.0% ✓", C_GREEN, True), "6,758"],
                [("GRU", C_GREEN, True), ("100.0% ✓", C_GREEN, True), "5,382"],
                [("RNN", C_GREEN, True), ("100.0% ✓", C_GREEN, True), "2,630"],
                [("LTC", C_RED, False), ("36.9% ✗", C_RED, False), "6,936"]],
               left=4.65, top=2.18, col_widths=[1.3, 1.5, 1.65],
               font_size=11, row_height=0.37)
    add_text(slide,
             "LTC 仅 36.9%（欠拟合）：\n"
             "规则信号 + CPU 时间限制 → 训练不足\n"
             "RNN 参数效率最高（2,630，100%）",
             4.7, 4.0, 4.4, 0.85, font_size=12, color=C_LIGHT)

    # CartPole
    add_rect(slide, 9.3, 1.6, 3.7, 2.8, fill_color=RGBColor(0x04, 0x22, 0x3A))
    add_text(slide, "实验七：CartPole RL", 9.4, 1.7, 3.5, 0.4, font_size=15, bold=True, color=C_ACCENT)
    make_table(slide, ["模型", "最终奖励"],
               [[("CfC", C_GREEN, True), ("24.0 ✓", C_GREEN, True)],
                ["LSTM", "16.0"]],
               left=9.35, top=2.18, col_widths=[1.5, 2.1],
               font_size=12, row_height=0.48)
    add_text(slide,
             "CfC 奖励优于 LSTM\n"
             "两者均未达到解决阈值（≥195）\n"
             "REINFORCE 高方差 → 结论需更多实验",
             9.4, 3.2, 3.5, 1.0, font_size=12, color=C_LIGHT)

    # Summary section
    add_divider(slide, 4.55, width_in=12.7, left=0.3)
    add_text(slide, "横向小结", 0.3, 4.7, 12.7, 0.42, font_size=17, bold=True, color=C_ACCENT)
    summary = [
        "金融预测：随机游走信号，三模型趋同，LNN 参数代价难以转化为增益（适合轻量 GRU）",
        "UCI HAR：规则传感器信号，CfC/LSTM/GRU/RNN 全达 100%，LTC 受训练效率拖累",
        "CartPole RL：CfC 策略网络初步展示连续控制潜力，需低方差算法（PPO）进一步验证",
    ]
    bullet_list(slide, summary, 0.3, 5.2, 12.7, 2.3, font_size=13, spacing=0.55)


def slide_12_core_findings(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "06  综合分析：液态网络的优势边界",
                 "Core Findings – When Does LNN Win?")

    # Left: LNN wins
    add_rect(slide, 0.3, 1.6, 6.0, 5.55, fill_color=RGBColor(0x02, 0x38, 0x20))
    add_text(slide, "✅  LNN 显著优势场景", 0.45, 1.7, 5.7, 0.48,
             font_size=18, bold=True, color=C_GREEN)

    lnn_wins = [
        ("◆", "刚性动力学建模",
         "下课铃人流激增等突发尖峰 → CfC 误差约为 LSTM 的 44%（实验三）"),
        ("◆", "不规则采样 / 传感器缺失",
         "时间常数 τ 内置缺失感知机制 → 优于 LSTM 的普通掩码处理"),
        ("◆", "最高拟合精度（洁净数据）",
         "RMSE 和 R² 均优于全部传统 RNN（实验四：R²=0.967）"),
        ("◆", "连续控制任务",
         "CartPole RL 中 CfC 期望奖励高于 LSTM（实验七）"),
    ]
    y = 2.32
    for sym, title, desc in lnn_wins:
        add_text(slide, sym, 0.4, y, 0.4, 0.38, font_size=14, bold=True, color=C_GREEN)
        add_text(slide, title, 0.8, y, 5.2, 0.38, font_size=14, bold=True, color=C_WHITE)
        add_text(slide, desc, 0.4, y + 0.42, 5.8, 0.65, font_size=12, color=C_LIGHT)
        y += 1.15

    # Right: RNN wins
    add_rect(slide, 6.7, 1.6, 6.2, 5.55, fill_color=RGBColor(0x38, 0x10, 0x10))
    add_text(slide, "⚠️  传统 RNN 表现更优场景", 6.85, 1.7, 5.9, 0.48,
             font_size=18, bold=True, color=C_RED)

    rnn_wins = [
        ("◆", "规则采样、无刚性跳变的合成数据",
         "Walker2d 合成数据：LSTM MSE 约为 LTC 的 23%（实验二）"),
        ("◆", "随机游走类金融序列",
         "NASDAQ 预测三模型趋同，LNN 高参数量无额外增益（实验五）"),
        ("◆", "推理延迟敏感的边缘部署",
         "LSTM 仅 0.261 ms vs LTC 19.791 ms，约快 76 倍"),
        ("◆", "计算资源严格受限场景",
         "GRU/LSTM 参数量少、训练速度快；LTC 训练耗时高出 10–50 倍"),
    ]
    y = 2.32
    for sym, title, desc in rnn_wins:
        add_text(slide, sym, 6.8, y, 0.4, 0.38, font_size=14, bold=True, color=C_RED)
        add_text(slide, title, 7.2, y, 5.5, 0.38, font_size=14, bold=True, color=C_WHITE)
        add_text(slide, desc, 6.8, y + 0.42, 6.0, 0.65, font_size=12, color=C_LIGHT)
        y += 1.15

    # Bottom summary
    add_rect(slide, 0.3, 7.12, 12.6, 0.25, fill_color=C_ACCENT)


def slide_13_cfc_vs_ltc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "06b  CfC vs LTC：液态网络内部对比",
                 "CfC vs LTC – Within LNN Family")

    headers = ["维度", "CfC", "LTC", "胜出"]
    rows = [
        ["推理延迟", "5.775 ms", "19.791 ms", ("CfC ✓", C_GREEN, True)],
        ["分类准确率（UCI HAR）", "100.0%", "36.9%", ("CfC ✓", C_GREEN, True)],
        ["鲁棒性 RMSE（洁净）", "0.02900", "0.03394", ("CfC ✓", C_GREEN, True)],
        ["R² 评分", "0.9669", "0.9547", ("CfC ✓", C_GREEN, True)],
        ["参数量（32 单元）", "21,025", "5,542", ("LTC ✓", C_GREEN, True)],
        ["理论可解释性", "近似解析", "原始 ODE", ("LTC ✓", C_GREEN, True)],
    ]
    make_table(slide, headers, rows,
               left=1.5, top=1.65,
               col_widths=[3.5, 2.5, 2.5, 2.0],
               font_size=14, row_height=0.55)

    conclusion = (
        "综合结论：CfC 在性能和效率两个维度均优于 LTC，"
        "是液态网络系列中更具工程实用性的选择；"
        "LTC 因更接近生物神经元 ODE 原始形式，"
        "在理论研究和可解释性分析中仍具重要价值。"
    )
    add_rect(slide, 1.0, 5.42, 11.0, 1.2, fill_color=RGBColor(0x00, 0x50, 0x70))
    add_text(slide, conclusion, 1.15, 5.52, 10.7, 1.0,
             font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_14_params(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "06c  参数效率分析", "Parameter Efficiency Analysis")

    headers = ["模型", "参数量", "相对 LSTM", "典型任务最优表现"]
    rows = [
        ["RNN", "~2,600", "0.56×", "UCI HAR：100% 准确率（最高参数效率）"],
        [("GRU", C_ACCENT2, False), "3,489", "0.75×", "金融预测：MSE 最优；高噪声鲁棒性最佳"],
        ["LSTM", "4,641", "1.00×（基准）", "Walker2d：MSE 最优；推理延迟最低（0.261 ms）"],
        ["LTC", "5,542", "1.19×", "理论基准；校园流量下优于 LSTM（实验三除外）"],
        [("CfC", C_GREEN, True), "21,025", ("4.53×", C_RED, False),
         "校园人流：MSE 仅 LSTM 的 44%；鲁棒性 R²=0.967 最优"],
    ]
    make_table(slide, headers, rows,
               left=0.4, top=1.65,
               col_widths=[1.5, 1.8, 2.0, 7.5],
               font_size=13, row_height=0.55)

    note = ("CfC 参数量约为 LSTM 的 4.5 倍，主要源于 ncps 库完整闭合形式参数矩阵。"
            "在性能显著优越的任务（如校园人流，误差减半）中这一代价合理；"
            "在性能相近的任务（如金融预测）中，较高参数量带来了不必要的计算和存储开销。")
    add_rect(slide, 0.4, 4.75, 12.5, 1.2, fill_color=RGBColor(0x04, 0x28, 0x42))
    add_text(slide, "💡 " + note, 0.55, 4.85, 12.2, 1.0, font_size=13, color=C_LIGHT)

    add_text(slide, "选型建议", 0.4, 6.1, 12.5, 0.42, font_size=17, bold=True, color=C_ACCENT)
    recs = [
        "刚性/缺失场景（校园人流、工业传感器）→ 首选 CfC",
        "规则信号、延迟敏感 → 首选 LSTM / GRU",
        "参数极度受限 → 首选 RNN 或 GRU",
        "理论研究 / 可解释性 → 首选 LTC",
    ]
    bullet_list(slide, recs, 0.4, 6.6, 12.5, 1.5, font_size=13, spacing=0.45, bullet="→")


def slide_15_contributions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "07  项目成果与贡献", "Contributions & Deliverables")

    categories = [
        ("💻 代码成果", [
            "统一模型库（src/models/registry.py）：LTC/CfC/LSTM/GRU/RNN 五种模型统一接口",
            "校园人流数据生成器（src/data/campus_flow.py）：可配置刚性尖峰、缺失率、噪声",
            "七个独立实验脚本（experiments/）：时序回归 ×4、分类、鲁棒性、强化学习",
            "一键运行入口（run_all.py）：完整模式 + 快速冒烟测试模式（--fast，5 轮）",
        ]),
        ("📊 实验成果", [
            "14 个 CSV 指标文件 + 15 幅可视化图表，全部保存至 results/ 子目录",
            "首次将 CfC 应用于同济大学校园人流预测场景，MSE 降低约 56%",
            "完整四维性能画像：准确性（MSE/RMSE/MAE/R²/Acc）、效率、鲁棒性、推理速度",
        ]),
        ("📄 文档成果", [
            "立项书（docs/proposal.md）：研究背景、目标、技术路线与实验设计",
            "超参数日志（docs/SITP_Log.md）：逐实验全部超参数配置，保障可复现性",
            "结题报告（docs/Final_Report_SITP.md）：学术规范组织的完整实验分析",
            "文献综述 + 学习笔记：为国内研究者提供中文液态神经网络学习资源",
        ]),
    ]

    x_positions = [0.3, 4.55, 8.8]
    for idx, (cat, items) in enumerate(categories):
        x = x_positions[idx]
        add_rect(slide, x, 1.6, 4.05, 5.5, fill_color=RGBColor(0x04, 0x26, 0x3E))
        add_rect(slide, x, 1.6, 4.05, 0.55, fill_color=C_HEADER)
        add_text(slide, cat, x + 0.1, 1.68, 3.85, 0.42, font_size=15, bold=True, color=C_ACCENT)
        y = 2.25
        for item in items:
            add_text(slide, "◆", x + 0.1, y, 0.35, 0.38, font_size=11, bold=True, color=C_ACCENT)
            add_text(slide, item, x + 0.45, y, 3.55, 0.8, font_size=11, color=C_LIGHT)
            y += 1.1


def slide_16_challenges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "08  挑战与未来展望", "Challenges & Future Work")

    # Left: Challenges
    add_rect(slide, 0.3, 1.6, 5.9, 5.55, fill_color=RGBColor(0x06, 0x20, 0x30))
    add_text(slide, "⚙️ 遇到的主要挑战", 0.45, 1.7, 5.6, 0.45, font_size=18, bold=True, color=C_RED)

    challenges = [
        ("①", "LTC 训练效率",
         "ODE 数值求解导致训练时间比传统 RNN 高出 10–50 倍，"
         "UCI HAR 等大样本任务中严重限制可用数据量（最大 2000 条）。"),
        ("②", "ncps 库接口限制",
         "FullyConnected wiring 对分类任务适配性不如 AutoNCP；"
         "CfC 参数量显著高于理论预期，部分源于库的内部实现细节。"),
        ("③", "强化学习高方差",
         "REINFORCE 梯度方差大，CartPole 实验 500 episodes 内"
         "两种模型均未稳定达到解决阈值，难以得出强结论。"),
        ("④", "真实数据获取困难",
         "校园人流和 Walker2d 均依赖合成数据，"
         "真实数据获取存在工程和权限障碍。"),
    ]
    y = 2.3
    for num, title, desc in challenges:
        add_text(slide, num, 0.4, y, 0.45, 0.38, font_size=15, bold=True, color=C_RED)
        add_text(slide, title, 0.85, y, 5.2, 0.38, font_size=14, bold=True, color=C_WHITE)
        add_text(slide, desc, 0.4, y + 0.42, 5.7, 0.75, font_size=12, color=C_LIGHT)
        y += 1.25

    # Right: Future Work
    add_rect(slide, 6.6, 1.6, 6.4, 5.55, fill_color=RGBColor(0x03, 0x28, 0x40))
    add_text(slide, "🚀 未来展望", 6.75, 1.7, 6.1, 0.45, font_size=18, bold=True, color=C_ACCENT)

    future = [
        ("→", "接入真实校园数据",
         "对接同济大学门禁刷卡或摄像头视频人流统计，"
         "在真实传感器噪声和缺失条件下进一步验证 CfC 优势。"),
        ("→", "进阶强化学习实验",
         "采用 PPO 等低方差算法，在 Walker2d 等连续控制环境"
         "中对比 LNN 与 LSTM 作为策略网络的性能。"),
        ("→", "NCP 稀疏架构探索",
         "引入 AutoNCP wiring，研究其在序列分类和控制任务"
         "上相对 FullyConnected LTC 的优势。"),
        ("→", "边缘部署优化",
         "探索 CfC 模型量化压缩（INT8/FP16），"
         "评估在树莓派等嵌入式设备上的实际部署可行性。"),
        ("→", "多变量时序扩展",
         "拓展到气象、交通等高维多变量时序场景，"
         "评估液态网络在复杂时序依赖结构中的泛化能力。"),
    ]
    y = 2.3
    for arrow, title, desc in future:
        add_text(slide, arrow, 6.7, y, 0.45, 0.38, font_size=15, bold=True, color=C_ACCENT)
        add_text(slide, title, 7.15, y, 5.7, 0.38, font_size=14, bold=True, color=C_WHITE)
        add_text(slide, desc, 6.7, y + 0.42, 6.2, 0.7, font_size=12, color=C_LIGHT)
        y += 1.15


def slide_17_refs(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    slide_header(slide, "09  参考文献", "References")

    refs = [
        ("[1]", "Hasani, R. et al. (2021). Liquid Time-constant Networks. AAAI 2021, 35(9), 7657–7666."),
        ("[2]", "Hasani, R. et al. (2022). Closed-form Continuous-time Neural Networks. Nature Machine Intelligence, 4, 992–1003."),
        ("[3]", "Lechner, M. et al. (2020). Neural Circuit Policies Enabling Auditable Autonomy. Nature Machine Intelligence, 2, 642–652."),
        ("[4]", "arXiv:2510.07578v1 — 液态神经网络最新进展预印本. https://arxiv.org/abs/2510.07578"),
        ("[5]", "Hochreiter, S. & Schmidhuber, J. (1997). Long Short-Term Memory. Neural Computation, 9(8), 1735–1780."),
        ("[6]", "Cho, K. et al. (2014). Learning Phrase Representations using RNN Encoder-Decoder. EMNLP 2014."),
        ("[7]", "Elman, J. L. (1990). Finding Structure in Time. Cognitive Science, 14(2), 179–211."),
        ("[8]", "Williams, R. J. (1992). Simple Statistical Gradient-Following Algorithms for Connectionist RL. Machine Learning, 8, 229–256."),
        ("[9]", "Anguita, D. et al. (2013). A Public Domain Dataset for Human Activity Recognition Using Smartphones. ESANN 2013."),
    ]

    y = 1.65
    for num, ref in refs:
        add_rect(slide, 0.4, y, 0.65, 0.45, fill_color=C_ACCENT)
        add_text(slide, num, 0.43, y + 0.05, 0.62, 0.38, font_size=12, bold=True,
                 color=C_BG, align=PP_ALIGN.CENTER)
        add_text(slide, ref, 1.15, y + 0.03, 11.7, 0.42, font_size=12, color=C_LIGHT)
        y += 0.55


def slide_18_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    add_rect(slide, 0, 0, 13.33, 0.12, fill_color=C_ACCENT)
    add_rect(slide, 0, 7.38, 13.33, 0.12, fill_color=C_ACCENT)
    add_rect(slide, 0, 0.12, 0.12, 7.26, fill_color=C_ACCENT2)

    add_rect(slide, 1.0, 1.0, 11.3, 5.8, fill_color=RGBColor(0x04, 0x26, 0x3E))

    add_text(slide, "感谢聆听", 1.2, 1.4, 11.0, 1.0,
             font_size=48, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_divider(slide, 2.65, width_in=10.0, left=1.65)

    add_text(slide, "Thank You", 1.2, 2.85, 11.0, 0.7,
             font_size=28, italic=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

    add_text(slide, "液态神经网络与普通神经网络的对比研究", 1.5, 3.65, 10.3, 0.55,
             font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "同济大学 SITP 项目 · 2025–2026", 1.5, 4.25, 10.3, 0.45,
             font_size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "团队成员：陈山璞  韩兆津", 1.5, 4.8, 10.3, 0.45,
             font_size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, 3.5, 5.55, 6.3, 0.6, fill_color=C_ACCENT)
    add_text(slide, "欢迎提问与交流 🎓", 3.6, 5.62, 6.1, 0.45,
             font_size=18, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    # GitHub link
    add_text(slide, "🔗 https://github.com/ZhaoJinHan518/liquid_neural_network",
             1.5, 6.4, 10.3, 0.4, font_size=13, color=C_ACCENT2, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  PDF EXPORT
# ═══════════════════════════════════════════════════════════════════════════════

def _find_libreoffice():
    """Return the LibreOffice executable path, or None if not found."""
    for candidate in ("libreoffice", "soffice"):
        try:
            result = subprocess.run(
                [candidate, "--version"],
                capture_output=True, timeout=10
            )
            if result.returncode == 0:
                return candidate
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    return None


def export_to_pdf(pptx_path: str) -> str:
    """
    Convert *pptx_path* to PDF using LibreOffice.

    Returns the output PDF path on success.
    Raises RuntimeError with installation instructions if LibreOffice is absent.
    """
    lo = _find_libreoffice()
    if lo is None:
        raise RuntimeError(
            "LibreOffice 未找到，无法自动导出 PDF。\n\n"
            "请选择以下任一方式安装 LibreOffice：\n"
            "  Ubuntu/Debian : sudo apt install libreoffice\n"
            "  macOS         : brew install --cask libreoffice\n"
            "  Windows       : https://www.libreoffice.org/download/\n\n"
            "安装后重新运行：python docs/generate_ppt.py --pdf\n\n"
            "或者在 Microsoft PowerPoint / WPS 演示 中直接打开 .pptx 文件，\n"
            "选择「文件 → 导出 → 导出为 PDF」即可。"
        )

    out_dir = os.path.dirname(os.path.abspath(pptx_path))
    cmd = [lo, "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_path]
    print(f"  正在调用 LibreOffice 转换 PDF…  ({' '.join(cmd)})")
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        raise RuntimeError(
            f"LibreOffice 转换失败（返回码 {result.returncode}）：\n"
            f"{result.stderr.strip()}"
        )
    # LibreOffice names the output file after the input stem
    pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
    return pdf_path


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def _parse_args():
    parser = argparse.ArgumentParser(
        description="生成 SITP 结题汇报 PPT（可选同时导出 PDF）",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "示例：\n"
            "  python docs/generate_ppt.py\n"
            "  python docs/generate_ppt.py --out /tmp/my_slides.pptx\n"
            "  python docs/generate_ppt.py --pdf\n"
            "  python docs/generate_ppt.py --pdf --out /tmp/sitp_slides.pptx\n"
        ),
    )
    default_out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "SITP_结题汇报.pptx")
    parser.add_argument(
        "--out", metavar="OUTPUT.pptx", default=default_out,
        help=f"PPTX 输出路径（默认：{default_out}）",
    )
    parser.add_argument(
        "--pdf", action="store_true",
        help="同时将 PPTX 转换为 PDF（需要安装 LibreOffice）",
    )
    return parser.parse_args()


def main():
    args = _parse_args()

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_01_cover(prs)          ; print("  [1/18] Cover")
    slide_02_agenda(prs)         ; print("  [2/18] Agenda")
    slide_03_background(prs)     ; print("  [3/18] Background")
    slide_04_goals(prs)          ; print("  [4/18] Goals & Innovations")
    slide_05_models(prs)         ; print("  [5/18] Models & Theory")
    slide_06_exp_overview(prs)   ; print("  [6/18] Experiment Overview")
    slide_07_exp1_sine(prs)      ; print("  [7/18] Exp1: Damped Sine")
    slide_08_exp2_walker(prs)    ; print("  [8/18] Exp2: Walker2d")
    slide_09_campus_flow(prs)    ; print("  [9/18] Exp3: Campus Flow")
    slide_10_robustness(prs)     ; print(" [10/18] Exp4: Robustness")
    slide_11_other_exps(prs)     ; print(" [11/18] Exp5-7: Finance/HAR/CartPole")
    slide_12_core_findings(prs)  ; print(" [12/18] Core Findings")
    slide_13_cfc_vs_ltc(prs)     ; print(" [13/18] CfC vs LTC")
    slide_14_params(prs)         ; print(" [14/18] Parameter Efficiency")
    slide_15_contributions(prs)  ; print(" [15/18] Contributions")
    slide_16_challenges(prs)     ; print(" [16/18] Challenges & Future")
    slide_17_refs(prs)           ; print(" [17/18] References")
    slide_18_thanks(prs)         ; print(" [18/18] Thank You")

    out_path = os.path.abspath(args.out)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    print(f"\n✅ PPTX 已保存：{out_path}")

    if args.pdf:
        try:
            pdf_path = export_to_pdf(out_path)
            print(f"✅ PDF  已保存：{pdf_path}")
        except RuntimeError as exc:
            print(f"\n⚠️  PDF 导出失败：\n{exc}", file=sys.stderr)
            sys.exit(1)


if __name__ == "__main__":
    main()
